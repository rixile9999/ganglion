#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Ganglion presentation slide generator using python-pptx.

Light theme (cool clean: white / slate / teal-amber) with a single running
IoT-light example threaded through every content slide so the audience can
see how the initial input value transforms at each pipeline stage.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ==============================================================================
# CONSTANTS & THEME DEFINITIONS  (light theme — cool clean)
# ==============================================================================
FONT_HEADER = "IBM Plex Mono"
FONT_BODY = "Arial"

# Background / surfaces
COLOR_INK         = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF  page background
COLOR_INK_CARD    = RGBColor(0xF8, 0xFA, 0xFC)   # #F8FAFC  card surface (slate-50)
COLOR_PANEL_BG    = RGBColor(0xF1, 0xF5, 0xF9)   # #F1F5F9  example-probe tint (slate-100)
COLOR_PANEL_SOFT  = RGBColor(0xEC, 0xFD, 0xF5)   # #ECFDF5  soft teal wash for hero box

# Ink / text
COLOR_BONE_BRIGHT = RGBColor(0x0F, 0x17, 0x2A)   # #0F172A  near-black (slate-900) — headings
COLOR_BONE        = RGBColor(0x33, 0x41, 0x55)   # #334155  body text (slate-700)
COLOR_BONE_DIM    = RGBColor(0x64, 0x74, 0x8B)   # #64748B  secondary (slate-500)
COLOR_BONE_MUTE   = RGBColor(0x94, 0xA3, 0xB8)   # #94A3B8  captions (slate-400)

# Rules
COLOR_RULE        = RGBColor(0xE2, 0xE8, 0xF0)   # #E2E8F0  soft rule (slate-200)
COLOR_RULE_STRONG = RGBColor(0xCB, 0xD5, 0xE1)   # #CBD5E1  strong rule (slate-300)

# Accents (WCAG-tuned for light bg)
COLOR_TEAL        = RGBColor(0x0F, 0x76, 0x6E)   # #0F766E  teal-700
COLOR_CHARTREUSE  = RGBColor(0x4D, 0x7C, 0x0F)   # #4D7C0F  lime-700
COLOR_AMBER       = RGBColor(0xB4, 0x53, 0x09)   # #B45309  amber-700
COLOR_VERMILLION  = RGBColor(0xB9, 0x1C, 0x1C)   # #B91C1C  red-700

# ==============================================================================
# RUNNING EXAMPLE  (threaded across every content slide)
# ==============================================================================
EXAMPLE_USER = "거실 조명을 50%로 켜줘"
EXAMPLE_IR_FULL = '{"calls":[{"tool":"set_light","args":{"room":"living","level":50}}]}'

# Speech notes dictionary
SPEECH_NOTES = {
    1: "교수님, 안녕하십니까. 오늘 발표할 연구 주제는 Ganglion 프로젝트로, 에지 디바이스 환경에서 경량 언어 모델이 어떻게 하면 고비용·고지연의 툴 스키마 오버헤드를 극복하고 대형 모델급의 정확도로 툴 호출(Tool Calling)을 수행할 수 있을지에 대한 시스템-모델 공동 설계 연구입니다.\n\n저희 프레임워크는 소형 모델의 학습 파이프라인인 Reflex-LM과 호스트 단의 검증/보정을 맡는 Myelin Compiler의 협력을 핵심 메커니즘으로 취합니다.\n\n발표 전반에 걸쳐 '거실 조명을 50%로 켜줘' 라는 동일한 IoT 사용자 입력이 각 모듈을 거치며 어떻게 변환되는지를 슬라이드 상단의 EXAMPLE PROBE로 따라가며 시연합니다.",
    2: "본 연구의 핵심 가설은 두 가지입니다. 첫째, 도구 명세 전체를 매번 프롬프트에 실어 나르는 대신, 압축된 중간 언어인 Action IR을 사용해 토큰을 극적으로 줄이면서도 정확도를 유지할 수 있는가?\n\n둘째, 소형 모델의 추론 한계를 모델 가중치 튜닝과 호스트 컴파일러 최적화의 하이브리드 결합으로 메울 수 있는가 하는 점입니다. 개발자가 기능 명세(Spec)만 던지면 학습 데이터 생성부터 최종 서빙 아티팩트까지 자동 생성하는 Spec-Driven Model Factory의 비전을 담고 있습니다.\n\n상단 EXAMPLE PROBE: 같은 입력 '거실 조명을 50%로 켜줘'에 대해 Native tool-call은 ~287 토큰의 schema를 요구하지만 Action IR은 28 토큰으로 동일한 의미를 전달합니다.",
    3: "저희 시스템 아키텍처는 보시는 바와 같이 9단계의 순환 루프(9-Stage Closed Loop)로 구현되어 있으며, 이를 생물학적 신경절 조직에 빗대어 세 가지 모듈로 분리했습니다.\n\n자극을 수용하고 변환하는 Dendrite(Contract), 실질적 추론 발화를 수행하는 Soma(LM), 그리고 실행 로그를 분석하고 오류를 정적/동적으로 교정하는 Glia(Analyzer) 레이어로 구성되어 모듈 간 결합도를 최소화했습니다.\n\n상단 EXAMPLE PROBE: 동일한 IoT 입력이 9단계 중 05·TRACE 단계에서 어떻게 정형 트레이스 레코드로 응축되는지 보여줍니다.",
    4: "첫 번째 모듈인 학습 페이즈(Learning Phase)입니다. 도구 명세가 들어오면 자동으로 다양한 학습 템플릿과 합성 데이터를 생성합니다.\n\n이후 극소량의 골드 데이터셋을 통해 초소형 모델(Reflex-LM)에 SFT를 적용하여 포맷 학습을 끝내고, 툴 간의 정확한 선택 선호도를 높이기 위해 DPO 정렬 파이프라인까지 연동 가능한 구조를 완비했습니다.\n\n상단 EXAMPLE PROBE: 본 IoT 입력이 SYNTH 단계에서 (prompt, gold_ir) 골드 페어로 변환되어 SFT 학습 row로 사용됩니다.",
    5: "두 번째 핵심 모듈인 컴파일러 패스(Compiler Passes)입니다. 저희는 모델의 생성을 그냥 방치하지 않고 컴파일 타임 최적화를 적용했습니다.\n\n첫째로, 도구 명세로부터 EBNF 문법을 실시간 생성하여 디코딩 시점에서 유효한 Action IR만 출력되도록 토큰 마스킹을 수행합니다. 둘째로, 그럼에도 발생하는 인자 누락이나 논리 에러는 호스트 단의 Myelin Compiler가 감지하여 디폴트 값을 동적으로 주입하거나, 에러 로그 피드백을 통해 런타임 재시도를 유도하는 이중 방어 체계를 갖췄습니다.\n\n상단 EXAMPLE PROBE: 모델이 비정규(room='거실') 값을 뱉고 level 인자를 누락해도 alias 매핑(거실→living)과 default 주입(level=100)으로 즉시 복원됩니다.",
    6: "세 번째 모듈인 셀프 부트스트랩 데이터 루프입니다. 에이전트의 수행 기록은 모두 Trace 데이터 모델로 수집됩니다. 이후 결정론적 분류기인 Taxonomy를 통해 오류가 파싱 에러인지, 인자 누락인지를 체계적으로 라벨링합니다.\n\n최종적으로 보정이나 재시도를 거쳐 '성공으로 복구된 Action IR' 로그를 정제하여 다음 세대 모델의 SFT 데이터로 재피딩(Re-feeding)함으로써 모델을 점점 강건하게 진화시키는 자가 증강 루프를 지원합니다.\n\n상단 EXAMPLE PROBE: 본 IoT 입력의 trace.repaired 레코드가 다음 세대 SFT 골드 row로 환원되는 과정을 보여줍니다.",
    7: "저희가 수행한 정량적 평가 결과입니다. 자체 IoT 도메인 평가에서는 50개 도구 확장 시 입력 토큰을 68.5% 절감하여 스케일링 효율을 증명했고, 평균 Latency를 19% 단축함과 동시에 속도의 편차도 줄였습니다.\n\n특히 Qwen 0.6B 초소형 모델 기준 Berkeley Function Calling Benchmark(BFCL v4) 평가에서, 기존 Native tool-calling baseline의 31.4% 대비 압도적인 91.2%의 AST Match 종합 정확도(+59.8%pp)를 기록하였습니다.\n\n또한 호출 거부 문맥(Irrelevance)에서도 No-Call 정확도 100%를 달성하였으며, 이는 명세 기반 SFT와 Myelin Compiler의 사후 규칙 보정(Post-Correction) 조합이 초소형 모델의 한계를 완전히 극복할 수 있음을 증명합니다.\n\n상단 EXAMPLE PROBE: 본 IoT 입력 단일 케이스에서 측정된 Native vs Ganglion 비교 수치입니다 (287→92 tok, 482→391 ms).",
    8: "마지막으로 향후 개발 및 연구 계획(TBD)입니다. 엔지니어링 측면에서는 모듈 리네이밍 마이그레이션과 함께 범용 도구 표준 프로토콜인 MCP 연동 테스트를 준비 중입니다.\n\n학술적으로는 Holdout 세트에서 보이는 오버피팅을 방지하기 위한 데이터 다양성 증강 기법을 연구하고, 저희가 확보한 강력한 Verifier 환경을 활용해 GRPO 기반의 강화학습을 Reflex-LM에 적용하여 추론 성능을 추가로 극대화할 생각입니다.\n\n이번 학기에 논문 초안 작성을 목표로 하고 있으며, 시스템적인 기여도를 부각해 MLSys 2026 투고를 생각하고 있는데, 교수님의 의견을 여쭙고 싶습니다. 감사합니다."
}

# ==============================================================================
# UTILITY HELPER FUNCTIONS
# ==============================================================================
def set_shape_card_bg(shape):
    """Light-card surface for content boxes (white-ish fill + soft border)."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_INK_CARD
    shape.line.color.rgb = COLOR_RULE_STRONG
    shape.line.width = Pt(1.0)


def apply_background(slide):
    """Solid white background for the slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_INK


def add_slide_header(slide, title_text, slide_num, tag_text=None):
    """Top header block: kicker tag, slide title, page number, divider rule."""
    if tag_text:
        tag_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.25), Inches(8.0), Inches(0.35))
        tf = tag_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = tag_text.upper()
        p.font.name = FONT_HEADER
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEAL

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(9.0), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_HEADER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_BONE_BRIGHT

    num_box = slide.shapes.add_textbox(Inches(10.5), Inches(0.55), Inches(2.0), Inches(0.7))
    tf = num_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = f"[ {slide_num:02d} / 08 ]"
    p.font.name = FONT_HEADER
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_BONE_MUTE

    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.75), Inches(1.3), Inches(12.58), Inches(1.3))
    line.line.color.rgb = COLOR_RULE
    line.line.width = Pt(1.0)


def add_slide_footer(slide):
    """Bottom divider rule + brand stamp + date."""
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.75), Inches(6.7), Inches(12.58), Inches(6.7))
    line.line.color.rgb = COLOR_RULE
    line.line.width = Pt(1.0)

    brand_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.85), Inches(6.0), Inches(0.4))
    tf = brand_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = "GANGLION // SYSTEMS CO-DESIGN"
    p.font.name = FONT_HEADER
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_BONE_MUTE

    date_box = slide.shapes.add_textbox(Inches(9.5), Inches(6.85), Inches(3.08), Inches(0.4))
    tf = date_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = "MAY 2026"
    p.font.name = FONT_HEADER
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_BONE_MUTE


def add_speaker_notes(slide, slide_idx):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = SPEECH_NOTES.get(slide_idx, "")


def add_example_panel(slide, stage_caption, transform_text, accent=None):
    """
    Compact 'EXAMPLE PROBE' bar (y=1.40 → 1.85) showing how the running
    IoT-light example transforms at the current pipeline stage.
    Shifts on-slide body content downward by ~0.35".
    """
    accent = accent or COLOR_TEAL

    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.75), Inches(1.40), Inches(11.83), Inches(0.45)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLOR_PANEL_BG
    panel.line.color.rgb = accent
    panel.line.width = Pt(0.75)

    tf = panel.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.04)

    # Top line: probe label + original user input
    p1 = tf.paragraphs[0]
    r_label = p1.add_run()
    r_label.text = f"EXAMPLE · {stage_caption}"
    r_label.font.name = FONT_HEADER
    r_label.font.size = Pt(8.5)
    r_label.font.bold = True
    r_label.font.color.rgb = accent

    r_input_lbl = p1.add_run()
    r_input_lbl.text = "    USER INPUT → "
    r_input_lbl.font.name = FONT_HEADER
    r_input_lbl.font.size = Pt(8.5)
    r_input_lbl.font.color.rgb = COLOR_BONE_MUTE

    r_input = p1.add_run()
    r_input.text = f"“{EXAMPLE_USER}”"
    r_input.font.name = FONT_HEADER
    r_input.font.size = Pt(8.5)
    r_input.font.italic = True
    r_input.font.color.rgb = COLOR_BONE_DIM

    # Bottom line: transformation at this stage
    p2 = tf.add_paragraph()
    p2.space_before = Pt(1)
    r_arrow = p2.add_run()
    r_arrow.text = "▸  "
    r_arrow.font.name = FONT_HEADER
    r_arrow.font.size = Pt(10)
    r_arrow.font.bold = True
    r_arrow.font.color.rgb = accent
    r_xform = p2.add_run()
    r_xform.text = transform_text
    r_xform.font.name = FONT_HEADER
    r_xform.font.size = Pt(10)
    r_xform.font.bold = True
    r_xform.font.color.rgb = COLOR_BONE_BRIGHT


# ==============================================================================
# SLIDE BUILDERS
# ==============================================================================

def build_slide_1(prs):
    """Slide 1: Title — introduces the running example."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)

    # Outer decorator frame
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5))
    border.fill.background()
    border.line.color.rgb = COLOR_RULE
    border.line.width = Pt(1)

    # Kicker tag
    tag_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.6), Inches(3.5), Inches(0.4))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = COLOR_INK_CARD
    tag_box.line.color.rgb = COLOR_TEAL
    tag_box.line.width = Pt(1)
    tf = tag_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "MLSYS / NEURIPS RESEARCH OUTLINE"
    p.font.name = FONT_HEADER
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEAL

    # Main title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)

    p = tf.paragraphs[0]
    p.text = "Ganglion: "
    p.font.name = FONT_HEADER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_BONE_BRIGHT

    run = p.add_run()
    run.text = "Spec-Driven Model Factory"
    run.font.name = FONT_HEADER
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = COLOR_CHARTREUSE

    run2 = p.add_run()
    run2.text = "\nfor Edge AI Tool Calling"
    run2.font.name = FONT_HEADER
    run2.font.size = Pt(36)
    run2.font.bold = True
    run2.font.color.rgb = COLOR_BONE_BRIGHT

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.9))
    tf = sub_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = "중간 표현인 Action IR과 검증기인 Myelin Compiler의 공동 설계를 통한\n경량 언어 모델(Reflex-LM)의 툴 호출 가속화 및 강건성 확보 프레임워크"
    p.font.name = FONT_BODY
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_BONE_DIM

    # Running example foreshadow box
    fore = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.05), Inches(10.3), Inches(0.85))
    fore.fill.solid()
    fore.fill.fore_color.rgb = COLOR_PANEL_BG
    fore.line.color.rgb = COLOR_TEAL
    fore.line.width = Pt(0.75)
    tf_f = fore.text_frame
    tf_f.word_wrap = True
    tf_f.vertical_anchor = MSO_ANCHOR.TOP
    tf_f.margin_left = Inches(0.20)
    tf_f.margin_right = Inches(0.20)
    tf_f.margin_top = Inches(0.10)
    p_f1 = tf_f.paragraphs[0]
    p_f1.alignment = PP_ALIGN.LEFT
    r_f1 = p_f1.add_run()
    r_f1.text = "EXAMPLE PROBE"
    r_f1.font.name = FONT_HEADER
    r_f1.font.size = Pt(9.5)
    r_f1.font.bold = True
    r_f1.font.color.rgb = COLOR_TEAL
    r_f2 = p_f1.add_run()
    r_f2.text = "   동일한 IoT 사용자 입력이 각 모듈에서 어떻게 변환되는지를 매 슬라이드 상단의 PROBE 패널에서 추적합니다"
    r_f2.font.name = FONT_HEADER
    r_f2.font.size = Pt(9.5)
    r_f2.font.color.rgb = COLOR_BONE_DIM

    p_f2 = tf_f.add_paragraph()
    p_f2.alignment = PP_ALIGN.LEFT
    p_f2.space_before = Pt(3)
    r_f3 = p_f2.add_run()
    r_f3.text = "▸  "
    r_f3.font.name = FONT_HEADER
    r_f3.font.size = Pt(13)
    r_f3.font.bold = True
    r_f3.font.color.rgb = COLOR_TEAL
    r_f4 = p_f2.add_run()
    r_f4.text = f"“{EXAMPLE_USER}”"
    r_f4.font.name = FONT_HEADER
    r_f4.font.size = Pt(13)
    r_f4.font.bold = True
    r_f4.font.color.rgb = COLOR_BONE_BRIGHT

    # Presenter meta
    meta_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.05), Inches(10.3), Inches(0.65))
    tf = meta_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = "발표자: [이름]      |      지도교수: [교수명] 교수님\nGANGLION // SYSTEMS CO-DESIGN"
    p.font.name = FONT_HEADER
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_BONE_MUTE

    add_speaker_notes(slide, 1)


def build_slide_2(prs):
    """Slide 2: Core Concepts — Action IR tokenization win."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    add_slide_header(slide, "연구 핵심 컨셉 및 가설", 2, "Core Concept")
    add_example_panel(
        slide,
        "STAGE · ACTION IR EMIT",
        'IR  {"calls":[{"tool":"set_light","args":{"room":"living","level":50}}]}   ·   28 tok  vs  Native schema 287 tok',
        accent=COLOR_CHARTREUSE,
    )

    cards_data = [
        {"x": 0.75, "y": 1.95, "w": 5.8, "h": 2.10, "title": "Action IR 중심의 발화",
         "desc": "모델이 장황한 API Schema 전체를 분석하고 출력하는 대신, 극도로 규격화되고 압축된 중간 표현(Action IR)만을 출력하게 제약하여 토큰 소모를 최소화함."},
        {"x": 6.78, "y": 1.95, "w": 5.8, "h": 2.10, "title": "Spec-Driven Model Factory",
         "desc": "도구 Catalog 명세 하나만으로 학습데이터 합성, SFT/DPO 모델 파인튜닝, 디코딩 EBNF 문법 규칙 제어 생성을 자동으로 구축하는 통합 팩토리 플랫폼."},
        {"x": 0.75, "y": 4.20, "w": 5.8, "h": 2.10, "title": "Tuning & Compilation 하이브리드",
         "desc": "경량화된 가중치 튜닝(Reflex-LM)과 호스트단의 정적/동적 검증 컴파일러 패스(Myelin Compiler)를 결합하여 소형 모델 추론 성능의 근본적 오버헤드를 보완."},
        {"x": 6.78, "y": 4.20, "w": 5.8, "h": 2.10, "title": "핵심 가치 극대화",
         "desc": "Edge 기기의 리소스 제약 속에서 추론 토큰 소모량의 비약적인 절감(68%↑), Latency 단축(19%↑), 그리고 컴파일 검증을 통한 AST 일치율의 강건함을 증명함."},
    ]

    for card in cards_data:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(card["x"]), Inches(card["y"]), Inches(card["w"]), Inches(card["h"]))
        set_shape_card_bg(box)

        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.22)

        p_title = tf.paragraphs[0]
        p_title.text = f"› {card['title']}"
        p_title.font.name = FONT_HEADER
        p_title.font.size = Pt(17)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_CHARTREUSE
        p_title.space_after = Pt(8)

        p_desc = tf.add_paragraph()
        p_desc.text = card["desc"]
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(12.5)
        p_desc.font.color.rgb = COLOR_BONE

    add_slide_footer(slide)
    add_speaker_notes(slide, 2)


def build_slide_3(prs):
    """Slide 3: 9-Stage closed loop (diagram + example trace probe)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    add_slide_header(slide, "9-Stage Closed-Loop Pipeline", 3, "System Pipeline")
    add_example_panel(
        slide,
        "STAGE 05 · TRACE RECORD",
        'Trace{ id:42, prompt:"거실…", ir:{calls:[set_light(living,50)]}, ok:true, latency_ms:391 }',
        accent=COLOR_TEAL,
    )

    # Nodes — pushed below the EXAMPLE PROBE bar (panel ends at y=1.85)
    nodes = {
        "01": {"x": 0.75, "y": 2.20, "w": 1.9, "h": 1.0, "title": "01 · SYNTH",     "sub": "lm.synth",        "color": COLOR_CHARTREUSE},
        "02": {"x": 3.20, "y": 2.20, "w": 1.9, "h": 1.0, "title": "02 · FINETUNE",  "sub": "lm.finetune",     "color": COLOR_CHARTREUSE},
        "03": {"x": 5.65, "y": 1.95, "w": 1.9, "h": 0.7, "title": "03 · BENCH IOT", "sub": "consumer.iot",    "color": COLOR_BONE_DIM},
        "04": {"x": 5.65, "y": 2.80, "w": 1.9, "h": 0.7, "title": "04 · BENCH BFCL","sub": "consumer.bfcl",   "color": COLOR_BONE_DIM},
        "05": {"x": 8.10, "y": 2.20, "w": 1.9, "h": 1.0, "title": "05 · TRACE",     "sub": "analyzer.trace",  "color": COLOR_TEAL},
        "06": {"x": 8.10, "y": 4.30, "w": 1.9, "h": 1.0, "title": "06 · TAXONOMY",  "sub": "analyzer.errors", "color": COLOR_TEAL},
        "07": {"x": 5.65, "y": 4.30, "w": 1.9, "h": 1.0, "title": "07 · METRICS",   "sub": "analyzer.stats",  "color": COLOR_TEAL},
        "08": {"x": 3.20, "y": 4.30, "w": 1.9, "h": 1.0, "title": "08 · RULES",     "sub": "feedback.patch",  "color": COLOR_AMBER, "dash": True},
        "09": {"x": 0.75, "y": 4.30, "w": 1.9, "h": 1.0, "title": "09 · PATCH",     "sub": "contract.update", "color": COLOR_CHARTREUSE},
    }

    for nd in nodes.values():
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(nd["x"]), Inches(nd["y"]), Inches(nd["w"]), Inches(nd["h"]))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_INK_CARD
        box.line.color.rgb = nd["color"]
        box.line.width = Pt(1.5)
        if nd.get("dash"):
            box.line.dash_style = 2

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)

        p_title = tf.paragraphs[0]
        p_title.text = nd["title"]
        p_title.font.name = FONT_HEADER
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_BONE_BRIGHT

        p_sub = tf.add_paragraph()
        p_sub.text = nd["sub"]
        p_sub.font.name = FONT_HEADER
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = COLOR_BONE_MUTE

    # Connectors — derived from new y-coordinates (row1 mid=2.70, 03 mid=2.30, 04 mid=3.15, row2 mid=4.80)
    connectors = [
        (Inches(2.65), Inches(2.70), Inches(3.20), Inches(2.70), COLOR_CHARTREUSE),  # 01 -> 02
        (Inches(5.10), Inches(2.70), Inches(5.65), Inches(2.30), COLOR_CHARTREUSE),  # 02 -> 03
        (Inches(5.10), Inches(2.70), Inches(5.65), Inches(3.15), COLOR_CHARTREUSE),  # 02 -> 04
        (Inches(7.55), Inches(2.30), Inches(8.10), Inches(2.70), COLOR_TEAL),        # 03 -> 05
        (Inches(7.55), Inches(3.15), Inches(8.10), Inches(2.70), COLOR_TEAL),        # 04 -> 05
        (Inches(9.05), Inches(3.20), Inches(9.05), Inches(4.30), COLOR_TEAL),        # 05 -> 06 (down)
        (Inches(8.10), Inches(4.80), Inches(7.55), Inches(4.80), COLOR_TEAL),        # 06 -> 07
        (Inches(5.65), Inches(4.80), Inches(5.10), Inches(4.80), COLOR_TEAL),        # 07 -> 08
        (Inches(3.20), Inches(4.80), Inches(2.65), Inches(4.80), COLOR_AMBER),       # 08 -> 09
        (Inches(1.70), Inches(4.30), Inches(1.70), Inches(3.20), COLOR_CHARTREUSE),  # 09 -> 01 (loop up)
    ]
    for x1, y1, x2, y2, color in connectors:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        line.line.color.rgb = color
        line.line.width = Pt(1.5)

    # Forward / Feedback explanation
    desc_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.55), Inches(11.83), Inches(1.0))
    tf = desc_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)

    p1 = tf.paragraphs[0]
    r1a = p1.add_run()
    r1a.text = "· Forward Pass: "
    r1a.font.name = FONT_HEADER
    r1a.font.bold = True
    r1a.font.size = Pt(12.5)
    r1a.font.color.rgb = COLOR_TEAL
    r1b = p1.add_run()
    r1b.text = "도구 스펙 Catalog 분석 → 데이터 합성(01) → 소형 모델 파인튜닝(02) → 다양한 벤치마크 평가(03,04) → Trace 수집(05) 단계로 순차 실행."
    r1b.font.name = FONT_BODY
    r1b.font.size = Pt(12.5)
    r1b.font.color.rgb = COLOR_BONE
    p1.space_after = Pt(5)

    p2 = tf.add_paragraph()
    r2a = p2.add_run()
    r2a.text = "· Feedback Loop: "
    r2a.font.name = FONT_HEADER
    r2a.font.bold = True
    r2a.font.size = Pt(12.5)
    r2a.font.color.rgb = COLOR_AMBER
    r2b = p2.add_run()
    r2b.text = "에러 로그 분류(06) → 성능 통계 분석(07) → Myelin Compiler 컴파일 패스 규칙 도출(08) → 툴 계약 패치(09)의 순환 복구 구조를 가짐."
    r2b.font.name = FONT_BODY
    r2b.font.size = Pt(12.5)
    r2b.font.color.rgb = COLOR_BONE

    add_slide_footer(slide)
    add_speaker_notes(slide, 3)


def build_slide_4(prs):
    """Slide 4: Learning phase — synth pair + SFT/DPO."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    add_slide_header(slide, "Module 1. Learning Phase (Reflex-LM 학습)", 4, "Learning Phase")
    add_example_panel(
        slide,
        "STAGE 01 · SYNTH → GOLD PAIR",
        'Pair{ prompt:"거실 조명을 50%로 켜줘",  gold:{calls:[{tool:"set_light",args:{room:"living",level:50}}]} }',
        accent=COLOR_CHARTREUSE,
    )

    bullets = [
        ("Spec-Driven Data Synthesis (명세 기반 데이터 합성)",
         "입력된 툴 명세(ToolSpec)를 파싱하여, SFT/DPO 학습에 필요한 다양한 가상 에이전트 질문 시나리오와 이에 매칭되는 정답 Action IR 문자열 골드 쌍을 템플릿 엔진을 통해 자동 생성."),
        ("Supervised Fine-Tuning (SFT)",
         "경량 모델(Qwen3-0.6B)의 표현 한계를 개선하기 위해 Action IR 포맷 학습에 최적화된 SFT 실행. 툴당 최소 80~100여 건 수준의 미세 튜닝만으로도 에지 환경에서의 포맷 붕괴 극복."),
        ("Direct Preference Optimization (DPO)",
         "단순 문법 학습을 넘어, 여러 복수 도구 간의 모호한 문맥이나 잘못된 선택을 제어하고 권한 외 입력을 안전하게 필터링(No-Call Contract)하기 위해 Preference Alignment 학습 연동 지원."),
    ]

    for idx, (title, desc) in enumerate(bullets):
        y_pos = 1.95 + (idx * 1.50)

        ind = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(y_pos), Inches(0.08), Inches(1.25))
        ind.fill.solid()
        ind.fill.fore_color.rgb = COLOR_CHARTREUSE
        ind.line.fill.background()

        box = slide.shapes.add_textbox(Inches(1.0), Inches(y_pos), Inches(11.58), Inches(1.25))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)

        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = FONT_HEADER
        p_title.font.size = Pt(17)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_BONE_BRIGHT
        p_title.space_after = Pt(6)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_BONE_DIM

    add_slide_footer(slide)
    add_speaker_notes(slide, 4)


def build_slide_5(prs):
    """Slide 5: Compiler passes — grammar mask + post-correction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    add_slide_header(slide, "Module 2. Compiler Passes (Myelin Compiler)", 5, "Compiler Passes")
    add_example_panel(
        slide,
        "STAGE · GRAMMAR MASK + POST-CORRECT",
        'raw {tool:"set_light", args:{room:"거실"}}   →   alias(거실→living) + default(level=100)   →   fixed IR',
        accent=COLOR_AMBER,
    )

    # Left card — Constrained Decoding
    col1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.95), Inches(5.8), Inches(4.40))
    set_shape_card_bg(col1)
    tf1 = col1.text_frame
    tf1.word_wrap = True
    tf1.vertical_anchor = MSO_ANCHOR.TOP
    tf1.margin_left = tf1.margin_right = Inches(0.3)
    tf1.margin_top = Inches(0.3)

    p = tf1.paragraphs[0]
    p.text = "1. Constrained Decoding (Grammar)"
    p.font.name = FONT_HEADER
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEAL
    p.space_after = Pt(18)

    p_b1 = tf1.add_paragraph()
    p_b1.text = "· EBNF 구문 제약 Masking"
    p_b1.font.name = FONT_HEADER
    p_b1.font.size = Pt(14.5)
    p_b1.font.bold = True
    p_b1.font.color.rgb = COLOR_BONE_BRIGHT
    p_b1.space_after = Pt(4)

    p_b1_desc = tf1.add_paragraph()
    p_b1_desc.text = "Catalog 구조 정의 및 툴 규칙에 맞춰 BNF/EBNF 제약 엔진을 동적으로 빌드하고 엔진 레벨(xgrammar)에 바인딩합니다."
    p_b1_desc.font.name = FONT_BODY
    p_b1_desc.font.size = Pt(13)
    p_b1_desc.font.color.rgb = COLOR_BONE_DIM
    p_b1_desc.space_after = Pt(16)

    p_b2 = tf1.add_paragraph()
    p_b2.text = "· 출력 안정성 100% 보장"
    p_b2.font.name = FONT_HEADER
    p_b2.font.size = Pt(14.5)
    p_b2.font.bold = True
    p_b2.font.color.rgb = COLOR_BONE_BRIGHT
    p_b2.space_after = Pt(4)

    p_b2_desc = tf1.add_paragraph()
    p_b2_desc.text = "모델 추론 시점에 문법 규칙에 벗어나는 부적절한 JSON/DSL 토큰 후보를 마스킹 차단하여 포맷 에러를 차단합니다."
    p_b2_desc.font.name = FONT_BODY
    p_b2_desc.font.size = Pt(13)
    p_b2_desc.font.color.rgb = COLOR_BONE_DIM

    # Right card — Dynamic Post-Correction
    col2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.78), Inches(1.95), Inches(5.8), Inches(4.40))
    set_shape_card_bg(col2)
    tf2 = col2.text_frame
    tf2.word_wrap = True
    tf2.vertical_anchor = MSO_ANCHOR.TOP
    tf2.margin_left = tf2.margin_right = Inches(0.3)
    tf2.margin_top = Inches(0.3)

    p = tf2.paragraphs[0]
    p.text = "2. Dynamic Post-Correction"
    p.font.name = FONT_HEADER
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_AMBER
    p.space_after = Pt(18)

    p_b3 = tf2.add_paragraph()
    p_b3.text = "· Default Value Injection"
    p_b3.font.name = FONT_HEADER
    p_b3.font.size = Pt(14.5)
    p_b3.font.bold = True
    p_b3.font.color.rgb = COLOR_BONE_BRIGHT
    p_b3.space_after = Pt(4)

    p_b3_desc = tf2.add_paragraph()
    p_b3_desc.text = "언어 모델이 간혹 생략하는 필수 아규먼트(Arguments)를 파서 단에서 Catalog 기본 규격에 맞춰 동적 복구 보완합니다."
    p_b3_desc.font.name = FONT_BODY
    p_b3_desc.font.size = Pt(13)
    p_b3_desc.font.color.rgb = COLOR_BONE_DIM
    p_b3_desc.space_after = Pt(16)

    p_b4 = tf2.add_paragraph()
    p_b4.text = "· 샌드박스 Repair Loop"
    p_b4.font.name = FONT_HEADER
    p_b4.font.size = Pt(14.5)
    p_b4.font.bold = True
    p_b4.font.color.rgb = COLOR_BONE_BRIGHT
    p_b4.space_after = Pt(4)

    p_b4_desc = tf2.add_paragraph()
    p_b4_desc.text = "런타임 에러나 논리 위반 발생 시 에러 피드백을 프롬프트 형태로 샌드박스 루프 내에 재피딩하여 가혹 조건에서도 자가 복구 완료합니다."
    p_b4_desc.font.name = FONT_BODY
    p_b4_desc.font.size = Pt(13)
    p_b4_desc.font.color.rgb = COLOR_BONE_DIM

    add_slide_footer(slide)
    add_speaker_notes(slide, 5)


def build_slide_6(prs):
    """Slide 6: Self-bootstrapping — trace.repaired → next-gen gold."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    add_slide_header(slide, "Module 3. Self-Bootstrapping Data Loop", 6, "Self-Bootstrapping")
    add_example_panel(
        slide,
        "STAGE · TRACE → NEXT-GEN GOLD",
        'trace.repaired{ ok:true, ir:{calls:[set_light(living,50)]} }   ⟶   next-gen SFT row',
        accent=COLOR_CHARTREUSE,
    )

    # Left column — Telemetry & Taxonomy text
    box_left = slide.shapes.add_textbox(Inches(0.75), Inches(1.95), Inches(5.8), Inches(4.40))
    tf1 = box_left.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = Inches(0)

    p_t1 = tf1.paragraphs[0]
    p_t1.text = "통일된 Telemetry (Trace) 수집"
    p_t1.font.name = FONT_HEADER
    p_t1.font.size = Pt(17)
    p_t1.font.bold = True
    p_t1.font.color.rgb = COLOR_BONE_BRIGHT
    p_t1.space_after = Pt(6)

    p_d1 = tf1.add_paragraph()
    p_d1.text = "에이전트 입력 질의, 출력 토큰, Latency, 디코딩 전략 매개변수 및 Myelin Compiler의 사후 복구 횟수 등을 단일 정형화 규격인 Trace 데이터 모델로 로깅합니다."
    p_d1.font.name = FONT_BODY
    p_d1.font.size = Pt(13.5)
    p_d1.font.color.rgb = COLOR_BONE_DIM
    p_d1.space_after = Pt(22)

    p_t2 = tf1.add_paragraph()
    p_t2.text = "14개 범주의 Error Taxonomy"
    p_t2.font.name = FONT_HEADER
    p_t2.font.size = Pt(17)
    p_t2.font.bold = True
    p_t2.font.color.rgb = COLOR_BONE_BRIGHT
    p_t2.space_after = Pt(6)

    p_d2 = tf1.add_paragraph()
    p_d2.text = "도구 호출 에러 건에 대하여, 구문 에러/타입 불일치/매개변수 누락/의미론적 오류 등 14개 범주의 Taxonomy 분류 필터를 적용하여 결정론적으로 이상 원인을 태깅합니다."
    p_d2.font.name = FONT_BODY
    p_d2.font.size = Pt(13.5)
    p_d2.font.color.rgb = COLOR_BONE_DIM

    # Right column — feedback loop hero card
    box_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.78), Inches(1.95), Inches(5.8), Inches(4.40))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = COLOR_PANEL_SOFT
    box_right.line.color.rgb = COLOR_CHARTREUSE
    box_right.line.width = Pt(1.5)

    tf2 = box_right.text_frame
    tf2.word_wrap = True
    tf2.vertical_anchor = MSO_ANCHOR.TOP
    tf2.margin_left = tf2.margin_right = Inches(0.3)
    tf2.margin_top = Inches(0.35)

    p_loop_title = tf2.paragraphs[0]
    p_loop_title.text = "🔄 Self-Data Augmentation Loop"
    p_loop_title.font.name = FONT_HEADER
    p_loop_title.font.size = Pt(18)
    p_loop_title.font.bold = True
    p_loop_title.font.color.rgb = COLOR_CHARTREUSE
    p_loop_title.space_after = Pt(18)

    p_loop_desc = tf2.add_paragraph()
    p_loop_desc.text = (
        "초기 추론 과정에서 포맷 붕괴나 호출 실수로 에러를 유발했으나, Myelin Compiler의 "
        "샌드박스 Repair Loop를 통해 자가 복원 및 실행에 성공한 궤적 로그들을 정제하여 수집합니다.\n\n"
        "이 '성공 복원'된 트레이스 데이터를 다음 세대의 SFT/DPO 골드 학습 데이터셋으로 자동 환원 "
        "피딩함으로써, 추가적인 레이블링 비용 없이 모델의 동작 강건성을 자가 강화시키는 사이클을 완성합니다."
    )
    p_loop_desc.font.name = FONT_BODY
    p_loop_desc.font.size = Pt(13.5)
    p_loop_desc.font.color.rgb = COLOR_BONE
    p_loop_desc.line_spacing = 1.3

    add_slide_footer(slide)
    add_speaker_notes(slide, 6)


def build_slide_7(prs):
    """Slide 7: Quantitative evaluation — measured on the running case."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    add_slide_header(slide, "정량적 실험 평가 및 성과", 7, "Evaluation Results")
    add_example_panel(
        slide,
        "STAGE · MEASURED ON THIS CASE",
        "Native baseline:  287 tok / 482 ms      Ganglion:  92 tok / 391 ms       Δ  −68% tok,  −19% lat",
        accent=COLOR_TEAL,
    )

    # Left column — table
    rows, cols = 5, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.75), Inches(1.95), Inches(6.2), Inches(4.20))
    table = table_shape.table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.0)

    table_data = [
        ["벤치마크",            "비교 항목",            "Native 대비 성과"],
        ["IoT (50 tools)",      "입력 토큰 절감률",     "68.5% 절감"],
        ["IoT (n=250)",         "추론 지연 (Latency)",  "19.0% (339ms) 단축"],
        ["BFCL v4 (0.6B)",      "종합 정확도 (AST)",    "91.2% vs 31.4% (+59.8%)"],
        ["BFCL (Irrelevance)",  "No-Call 정확도",       "100.0% vs 0.0%"],
    ]

    for r_idx in range(rows):
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            cell.text = table_data[r_idx][c_idx]

            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_HEADER
            p.font.size = Pt(11) if r_idx == 0 else Pt(11.5)

            highlight = (c_idx == 2 and any(tok in table_data[r_idx][c_idx] for tok in ("절감", "단축", "+59.8%", "100.0%")))
            p.font.bold = (r_idx == 0) or highlight

            if r_idx == 0:
                p.font.color.rgb = COLOR_BONE_BRIGHT
            elif highlight:
                p.font.color.rgb = COLOR_CHARTREUSE
            else:
                p.font.color.rgb = COLOR_BONE

            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_PANEL_BG if r_idx == 0 else COLOR_INK_CARD

    # Right column — Capacity Cliff cards
    cliff_title = slide.shapes.add_textbox(Inches(7.3), Inches(1.95), Inches(5.3), Inches(0.5))
    tf_ct = cliff_title.text_frame
    tf_ct.word_wrap = True
    tf_ct.margin_left = tf_ct.margin_right = tf_ct.margin_top = tf_ct.margin_bottom = Inches(0)
    p_ct = tf_ct.paragraphs[0]
    p_ct.text = "📉 Reflex-LM (0.6B) Capacity Cliff 극복"
    p_ct.font.name = FONT_HEADER
    p_ct.font.size = Pt(16)
    p_ct.font.bold = True
    p_ct.font.color.rgb = COLOR_TEAL

    # Card 1 — Untuned baseline
    card1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(2.55), Inches(2.5), Inches(1.45))
    set_shape_card_bg(card1)
    card1.line.color.rgb = COLOR_AMBER
    card1.line.width = Pt(1.5)
    tf_c1 = card1.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_left = tf_c1.margin_right = Inches(0.15)
    tf_c1.margin_top = Inches(0.15)

    p_val1 = tf_c1.paragraphs[0]
    p_val1.text = "31.4%"
    p_val1.font.name = FONT_HEADER
    p_val1.font.size = Pt(28)
    p_val1.font.bold = True
    p_val1.font.color.rgb = COLOR_AMBER
    p_val1.alignment = PP_ALIGN.CENTER

    p_lbl1 = tf_c1.add_paragraph()
    p_lbl1.text = "Untuned Native Baseline\n(병렬 호출 및 거부 실패)"
    p_lbl1.font.name = FONT_BODY
    p_lbl1.font.size = Pt(10)
    p_lbl1.font.color.rgb = COLOR_BONE_DIM
    p_lbl1.alignment = PP_ALIGN.CENTER
    p_lbl1.space_before = Pt(4)

    # Card 2 — Tuned + Compiled
    card2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.1), Inches(2.55), Inches(2.5), Inches(1.45))
    set_shape_card_bg(card2)
    card2.line.color.rgb = COLOR_CHARTREUSE
    card2.line.width = Pt(1.5)
    tf_c2 = card2.text_frame
    tf_c2.word_wrap = True
    tf_c2.margin_left = tf_c2.margin_right = Inches(0.15)
    tf_c2.margin_top = Inches(0.15)

    p_val2 = tf_c2.paragraphs[0]
    p_val2.text = "91.2%"
    p_val2.font.name = FONT_HEADER
    p_val2.font.size = Pt(28)
    p_val2.font.bold = True
    p_val2.font.color.rgb = COLOR_CHARTREUSE
    p_val2.alignment = PP_ALIGN.CENTER

    p_lbl2 = tf_c2.add_paragraph()
    p_lbl2.text = "SFT + Myelin Compiler\n(사후 최적화 보정보완 완료)"
    p_lbl2.font.name = FONT_BODY
    p_lbl2.font.size = Pt(10)
    p_lbl2.font.color.rgb = COLOR_BONE_DIM
    p_lbl2.alignment = PP_ALIGN.CENTER
    p_lbl2.space_before = Pt(4)

    # Cliff explanation
    cliff_desc = slide.shapes.add_textbox(Inches(7.3), Inches(4.20), Inches(5.3), Inches(2.0))
    tf_cd = cliff_desc.text_frame
    tf_cd.word_wrap = True
    tf_cd.margin_left = tf_cd.margin_right = tf_cd.margin_top = tf_cd.margin_bottom = Inches(0)
    p_cd = tf_cd.paragraphs[0]
    p_cd.text = (
        "· 0.6B 초경량 파라미터 모델 특성상 학습 전에는 복수/병렬 도구 선택(Parallel Call) 시 "
        "심각한 문법 붕괴를 보이며 성공률이 0%에 수렴했었습니다.\n\n"
        "· 하지만 명세 기반 SFT와 Myelin Compiler의 EBNF 제약 및 동적 인자 보정(Post-Correction)을 "
        "융합 적용한 결과, 종합 Macro Avg 91.2%의 고성능 툴 호출을 달성하며 소형 에이전트 적용 가능성을 입증했습니다."
    )
    p_cd.font.name = FONT_BODY
    p_cd.font.size = Pt(11.5)
    p_cd.font.color.rgb = COLOR_BONE
    p_cd.line_spacing = 1.3

    add_slide_footer(slide)
    add_speaker_notes(slide, 7)


def build_slide_8(prs):
    """Slide 8: Roadmap & discussion (no example probe — forward-looking)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    add_slide_header(slide, "TBD & 향후 연구 계획 (Discussion)", 8, "Roadmap & Discussion")

    # Left column — Engineering TBDs
    col1 = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.8), Inches(4.7))
    tf1 = col1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = Inches(0)

    p_eng = tf1.paragraphs[0]
    p_eng.text = "🛠️ Engineering TBDs"
    p_eng.font.name = FONT_HEADER
    p_eng.font.size = Pt(18)
    p_eng.font.bold = True
    p_eng.font.color.rgb = COLOR_TEAL
    p_eng.space_after = Pt(20)

    p_e1 = tf1.add_paragraph()
    p_e1.text = "· 네임스페이스 Refactoring"
    p_e1.font.name = FONT_HEADER
    p_e1.font.size = Pt(14.5)
    p_e1.font.bold = True
    p_e1.font.color.rgb = COLOR_BONE_BRIGHT
    p_e1.space_after = Pt(4)
    p_e1_desc = tf1.add_paragraph()
    p_e1_desc.text = "레거시 코드를 말끔하게 제거하고 생체/신경학 구조를 본뜬 모듈 명세인 dendrite(규격), soma(추론), glia(분석) 테마로 최종 포팅 작업을 전개합니다."
    p_e1_desc.font.name = FONT_BODY
    p_e1_desc.font.size = Pt(13)
    p_e1_desc.font.color.rgb = COLOR_BONE_DIM
    p_e1_desc.space_after = Pt(16)

    p_e2 = tf1.add_paragraph()
    p_e2.text = "· MCP (Model Context Protocol) 연동"
    p_e2.font.name = FONT_HEADER
    p_e2.font.size = Pt(14.5)
    p_e2.font.bold = True
    p_e2.font.color.rgb = COLOR_BONE_BRIGHT
    p_e2.space_after = Pt(4)
    p_e2_desc = tf1.add_paragraph()
    p_e2_desc.text = "앤트로픽 주도의 범용 도구 표준 프로토콜 스펙을 지원하여, Ganglion 공정을 통한 범용 에이전트 인프라와의 결합 실효성을 테스트합니다."
    p_e2_desc.font.name = FONT_BODY
    p_e2_desc.font.size = Pt(13)
    p_e2_desc.font.color.rgb = COLOR_BONE_DIM

    # Right column — Academic Discussions
    col2 = slide.shapes.add_textbox(Inches(6.78), Inches(1.6), Inches(5.8), Inches(4.7))
    tf2 = col2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = Inches(0)

    p_acad = tf2.paragraphs[0]
    p_acad.text = "🎓 Academic Discussions"
    p_acad.font.name = FONT_HEADER
    p_acad.font.size = Pt(18)
    p_acad.font.bold = True
    p_acad.font.color.rgb = COLOR_AMBER
    p_acad.space_after = Pt(20)

    p_a1 = tf2.add_paragraph()
    p_a1.text = "· Holdout 일반화 및 오버피팅 극복"
    p_a1.font.name = FONT_HEADER
    p_a1.font.size = Pt(14.5)
    p_a1.font.bold = True
    p_a1.font.color.rgb = COLOR_BONE_BRIGHT
    p_a1.space_after = Pt(4)
    p_a1_desc = tf2.add_paragraph()
    p_a1_desc.text = "합성 데이터에 의존한 학습 시 발생하는 미학습 도메인 거동 저하를 막기 위해 Entity Swap, Context Variation 증강 가이드라인을 수립합니다."
    p_a1_desc.font.name = FONT_BODY
    p_a1_desc.font.size = Pt(13)
    p_a1_desc.font.color.rgb = COLOR_BONE_DIM
    p_a1_desc.space_after = Pt(16)

    p_a2 = tf2.add_paragraph()
    p_a2.text = "· Verifier-Driven GRPO 적용"
    p_a2.font.name = FONT_HEADER
    p_a2.font.size = Pt(14.5)
    p_a2.font.bold = True
    p_a2.font.color.rgb = COLOR_BONE_BRIGHT
    p_a2.space_after = Pt(4)
    p_a2_desc = tf2.add_paragraph()
    p_a2_desc.text = "Myelin Compiler의 정교한 룰 엔진을 무비용 보상 함수(Reward Function)로 활용하여, 강화학습(GRPO) 기반의 에이전트 행동 정렬 타당성을 자문합니다."
    p_a2_desc.font.name = FONT_BODY
    p_a2_desc.font.size = Pt(13)
    p_a2_desc.font.color.rgb = COLOR_BONE_DIM

    add_slide_footer(slide)
    add_speaker_notes(slide, 8)


# ==============================================================================
# MAIN
# ==============================================================================
def main():
    print("Initializing PPTX slide builder...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        ("Title (+ running example foreshadow)",       build_slide_1),
        ("Core Concepts (Action IR token win)",        build_slide_2),
        ("System Pipeline (trace probe)",              build_slide_3),
        ("Learning Phase (synth gold pair)",           build_slide_4),
        ("Compiler Passes (grammar + post-correct)",   build_slide_5),
        ("Self-Bootstrapping (trace → next-gen gold)", build_slide_6),
        ("Evaluation (measured on this case)",         build_slide_7),
        ("Future Work & Discussions",                  build_slide_8),
    ]
    for idx, (label, fn) in enumerate(builders, start=1):
        print(f"  Slide {idx}: {label}")
        fn(prs)

    output_dir = "web"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    output_path = os.path.join(output_dir, "slides.pptx")
    print(f"Saving compiled presentation to {output_path}...")
    prs.save(output_path)
    print("Slide deck creation successfully finished!")


if __name__ == "__main__":
    main()
